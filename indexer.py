# indexer.py
import os
import psutil
import faiss
import pickle
import numpy as np
from sentence_transformers import SentenceTransformer

# Import handlers for different file types
import PyPDF2
import docx
import openpyxl
from pptx import Presentation

# --- UPDATED: A much more capable text extraction function ---
def extract_text_from_file(file_path):
    """Extracts text content from a wide variety of common file types."""
    _, extension = os.path.splitext(file_path)
    extension = extension.lower()
    
    # A set of all extensions that can be treated as plain text
    plain_text_extensions = {
        '.txt', '.md', '.py', '.js', '.html', '.css', 
        '.json', '.xml', '.csv', '.log', '.ini', '.cfg'
    }

    try:
        if extension == '.pdf':
            with open(file_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                return "".join(page.extract_text() for page in reader.pages)
        
        elif extension == '.docx':
            doc = docx.Document(file_path)
            return "\n".join([para.text for para in doc.paragraphs])
            
        elif extension == '.pptx':
            prs = Presentation(file_path)
            text_runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(shape.text)
            return "\n".join(text_runs)
            
        elif extension == '.xlsx':
            workbook = openpyxl.load_workbook(file_path, read_only=True)
            text_runs = []
            for sheet in workbook:
                for row in sheet.iter_rows():
                    for cell in row:
                        if cell.value:
                            text_runs.append(str(cell.value))
            return "\n".join(text_runs)
            
        elif extension in plain_text_extensions:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()
                
    except Exception as e:
        # Silently ignore errors from corrupted or protected files
        return None
        
    # If the file extension is not in our supported list, return None
    return None

def build_index():
    """Crawls drives, attempts to process all files, and saves the index."""
    print("Initializing embedding model...")
    model = SentenceTransformer('all-MiniLM-L6-v2')

    blacklist = {
        'node_modules', 'venv', 'Lib', 'site-packages', '.git', '.cache',
        'AppData', 'Program Files', 'ProgramData', 'Windows', '$RECYCLE.BIN',
        '__pycache__', 'env', '.gradle', 'Program Files (x86)', 'MinGW', '.vscode'
    }

    all_embeddings = []
    file_path_map = []
    
    drives = [partition.mountpoint for partition in psutil.disk_partitions()]
    
    print(f"Starting file crawl across drives: {drives}")
    for drive in drives:
        for root, dirs, files in os.walk(drive):
            dirs[:] = [d for d in dirs if d not in blacklist]
            
            # --- CHANGE: Loop through ALL files, not just specific types ---
            for file in files:
                file_path = os.path.join(root, file)
                
                if any(f"{os.path.sep}{folder_name}{os.path.sep}" in file_path for folder_name in blacklist):
                    continue

                # The smart extractor function will now do the filtering
                content = extract_text_from_file(file_path)
                
                # This check naturally skips unsupported files (where content is None)
                if content and content.strip():
                    print(f"-> Indexing content from: {file_path}")
                    embedding = model.encode(content, convert_to_tensor=False)
                    all_embeddings.append(embedding)
                    file_path_map.append(file_path)

    if not all_embeddings:
        print("No files were found to index.")
        return

    # The rest of the function remains the same
    embeddings_np = np.array(all_embeddings).astype('float32')
    
    print("Building FAISS index...")
    index = faiss.IndexFlatL2(embeddings_np.shape[1])
    index.add(embeddings_np)
    
    script_dir = os.path.dirname(os.path.abspath(__file__))
    index_dir = os.path.join(script_dir, "index")
    os.makedirs(index_dir, exist_ok=True)
    faiss.write_index(index, os.path.join(index_dir, "file_content.index"))
    with open(os.path.join(index_dir, "file_path_map.pkl"), "wb") as f:
        pickle.dump(file_path_map, f)
        
    print(f"Indexing complete. Processed {len(file_path_map)} files.")

if __name__ == "__main__":
    build_index()